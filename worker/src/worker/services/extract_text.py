import subprocess
from pathlib import Path

import docx
import pandas as pd
import pdfplumber
from bs4 import BeautifulSoup
from pptx import Presentation


def extract_text(filepath: str, mime_type: str) -> str:
    if mime_type in {"text/plain", "text/markdown"}:
        return Path(filepath).read_text(encoding="utf-8")

    elif mime_type in {"text/html", "application/xml"}:
        with open(filepath, encoding="utf-8") as f:
            return BeautifulSoup(f, "lxml").get_text()

    elif mime_type == "text/csv":
        df = pd.read_csv(filepath)
        return df.to_string()

    elif mime_type == "application/pdf":
        with pdfplumber.open(filepath) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages)

    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return "\n".join(p.text for p in docx.Document(filepath).paragraphs)

    elif mime_type in {
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }:
        df = pd.read_excel(filepath)
        return df.to_string()

    elif mime_type in {
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }:
        prs = Presentation(filepath)
        return "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )

    elif mime_type in {
        "application/vnd.oasis.opendocument.text",
        "application/vnd.oasis.opendocument.spreadsheet",
        "application/vnd.oasis.opendocument.presentation",
    }:
        txt_path = filepath + ".txt"
        subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "txt:Text",
                filepath,
                "--outdir",
                Path(filepath).parent,
            ],
            check=True,
        )
        return Path(txt_path).read_text(encoding="utf-8")

    else:
        raise ValueError(f"Unsupported MIME type: {mime_type}")
